#!/usr/bin/env python3
"""Smoke-test GDPVal judge MCP servers.

Run inside the GDPVal base image, or anywhere the MCP CLIs and Python
dependencies are installed:

    python runs/smoke_gdpval_mcp.py
"""

from __future__ import annotations

import argparse
import json
import sys
import tempfile
from pathlib import Path
from typing import Any

from simple_agent_lab.evals.suites.gdpval.judge_mcp import (
    GDPVAL_MCP_READ_TOOL_NAMES,
    gdpval_mcp_server_configs,
    open_gdpval_judge_tools,
)
from simple_agent_lab.mcp import MCPConnection, connect_mcp


def main() -> None:
    parser = argparse.ArgumentParser(description="Smoke-test GDPVal MCP servers.")
    parser.add_argument("--workdir", default=None)
    parser.add_argument("--reference-dir", default=None)
    parser.add_argument("--keep-temp", action="store_true")
    args = parser.parse_args()

    temp = None
    if args.workdir and args.reference_dir:
        workdir = Path(args.workdir).resolve()
        reference_dir = Path(args.reference_dir).resolve()
        workdir.mkdir(parents=True, exist_ok=True)
        reference_dir.mkdir(parents=True, exist_ok=True)
    else:
        temp = tempfile.TemporaryDirectory()
        root = Path(temp.name)
        workdir = root / "workdir"
        reference_dir = root / "judge_inputs"
        workdir.mkdir()
        reference_dir.mkdir()

    try:
        samples = _write_samples(reference_dir)
        outside = reference_dir.parent / "outside.txt"
        outside.write_text("outside allowed roots", encoding="utf-8")
        failures: list[str] = []
        for config in gdpval_mcp_server_configs(
            workdir=workdir,
            reference_dir=reference_dir,
        ):
            try:
                _smoke_server(config, samples=samples, outside_path=outside)
            except Exception as exc:  # noqa: BLE001 - report every server cleanly
                failures.append(f"{config.name}: {type(exc).__name__}: {exc}")
        try:
            _smoke_filtered_agent_tools(workdir=workdir, reference_dir=reference_dir)
        except Exception as exc:  # noqa: BLE001 - report filter failures cleanly
            failures.append(f"agent-tool-filter: {type(exc).__name__}: {exc}")
        if failures:
            print("\nFAILED")
            for failure in failures:
                print(f"- {failure}")
            raise SystemExit(1)
        print("\nGDPVal MCP smoke ok")
    finally:
        if temp is not None and not args.keep_temp:
            temp.cleanup()


def _smoke_server(config: Any, *, samples: dict[str, Path], outside_path: Path) -> None:
    print(f"\n==> {config.name}: {config.command} {' '.join(config.args)}")
    with connect_mcp(config) as connection:
        tool_names = {tool.name for tool in connection.tools}
        if not tool_names:
            raise RuntimeError("server returned no tools")
        print(f"tools: {len(tool_names)}")
        _call_sample(
            connection, config.name, samples=samples, outside_path=outside_path
        )


def _call_sample(
    connection: MCPConnection,
    server_name: str,
    *,
    samples: dict[str, Path],
    outside_path: Path,
) -> None:
    if server_name == "filesystem":
        _assert_success(
            connection.call(
                "read_text_file",
                {"path": str(samples["txt"])},
                timeout=connection_tool_timeout(connection),
            )
        )
        denied = connection.call(
            "read_text_file",
            {"path": str(outside_path)},
            timeout=connection_tool_timeout(connection),
        )
        if not denied.isError:
            raise RuntimeError("filesystem MCP read outside allowed roots succeeded")
        print("sample read: txt; outside-root read denied")
        return

    if server_name == "pdf":
        result = connection.call(
            "read_pdf",
            {
                "sources": [{"path": str(samples["pdf"])}],
                "include_full_text": True,
            },
            timeout=connection_tool_timeout(connection),
        )
        _assert_success(result, contains="hello pdf")
        print("sample read: pdf")
        return

    if server_name == "excel":
        result = connection.call(
            "read_data_from_excel",
            {
                "filepath": str(samples["xlsx"]),
                "sheet_name": "Sheet1",
                "start_cell": "A1",
                "end_cell": "A1",
            },
            timeout=connection_tool_timeout(connection),
        )
        _assert_success(result, contains="alpha")
        print("sample read: xlsx")
        return

    if server_name == "word":
        result = connection.call(
            "get_document_text",
            {"filename": str(samples["docx"])},
            timeout=connection_tool_timeout(connection),
        )
        _assert_success(result, contains="hello word")
        print("sample read: docx")
        return

    if server_name == "ppt":
        opened = connection.call(
            "open_presentation",
            {"file_path": str(samples["pptx"])},
            timeout=connection_tool_timeout(connection),
        )
        _assert_success(opened, contains="presentation_id")
        payload = _json_from_result(opened)
        presentation_id = payload.get("presentation_id")
        if presentation_id:
            text = connection.call(
                "extract_presentation_text",
                {"presentation_id": presentation_id},
                timeout=connection_tool_timeout(connection),
            )
            _assert_success(text, contains="hello ppt")
        print("sample read: pptx")
        return

    raise RuntimeError(f"no sample call registered for server {server_name!r}")


def _smoke_filtered_agent_tools(*, workdir: Path, reference_dir: Path) -> None:
    expected = {
        f"{server}_{tool_name}"
        for server, tool_names in GDPVAL_MCP_READ_TOOL_NAMES.items()
        for tool_name in tool_names
    }
    with open_gdpval_judge_tools(
        workdir=workdir,
        reference_dir=reference_dir,
        mode="mcp",
    ) as tools:
        names = {tool.name for tool in tools}
    unexpected = names - expected
    if unexpected:
        raise RuntimeError(
            "MCP agent tool filter exposed non-allowlisted tools: "
            + ", ".join(sorted(unexpected))
        )
    if not names:
        raise RuntimeError("MCP agent tool filter exposed no tools")
    print(f"\nfiltered agent tools: {len(names)}")


def connection_tool_timeout(connection: MCPConnection) -> float:
    del connection
    return 90.0


def _assert_success(result: Any, *, contains: str | None = None) -> None:
    text = _text_from_result(result)
    if result.isError:
        raise RuntimeError(text or "tool returned isError")
    if contains and contains not in text:
        raise RuntimeError(f"tool result did not contain {contains!r}: {text[:500]}")


def _json_from_result(result: Any) -> dict[str, Any]:
    text = _text_from_result(result)
    try:
        payload = json.loads(text)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"expected JSON result, got: {text[:500]}") from exc
    if not isinstance(payload, dict):
        raise RuntimeError(
            f"expected JSON object result, got: {type(payload).__name__}"
        )
    return payload


def _text_from_result(result: Any) -> str:
    return "\n".join(str(getattr(item, "text", "")) for item in result.content)


def _write_samples(root: Path) -> dict[str, Path]:
    root.mkdir(parents=True, exist_ok=True)
    txt = root / "note.txt"
    txt.write_text("hello filesystem\n", encoding="utf-8")

    xlsx = root / "sample.xlsx"
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet["A1"] = "alpha"
    workbook.save(xlsx)

    docx = root / "sample.docx"
    from docx import Document

    document = Document()
    document.add_paragraph("hello word")
    document.save(docx)

    pptx = root / "sample.pptx"
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "hello ppt"
    presentation.save(pptx)

    pdf = root / "sample.pdf"
    from reportlab.pdfgen import canvas

    page = canvas.Canvas(str(pdf))
    page.drawString(72, 720, "hello pdf")
    page.save()

    return {
        "txt": txt,
        "xlsx": xlsx,
        "docx": docx,
        "pptx": pptx,
        "pdf": pdf,
    }


if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        sys.exit(130)
